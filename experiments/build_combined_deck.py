# -*- coding: utf-8 -*-
"""CB-SAFE single combined deck: 22 numbered content slides.
1-15  = instructor 'method presentation' structure (stops cleanly at 15).
16-22 = full-paper research walkthrough (results + figures).
DNA: 16:9, Segoe UI, blue kicker, no accent bar, muted footers, no em-dash.
Full-form per-slide citations; selective bold key-phrase on 6 slides."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

BLUE = RGBColor(0x2A, 0x78, 0xD6)
INK  = RGBColor(0x0B, 0x0B, 0x0B)
INK2 = RGBColor(0x26, 0x25, 0x1F)
MUT  = RGBColor(0x89, 0x87, 0x81)
CALL = RGBColor(0xF2, 0xF5, 0xFA)
BODY = "Segoe UI"
FIG  = r"d:/ISM_Quantum/results/figs/"

P = Presentation()
P.slide_width = Inches(13.333)
P.slide_height = Inches(7.5)
BLANK = P.slide_layouts[6]
TOTAL = 22

# ---- full-form references (exact vol/no/pp from the paper bibliography) ----
REF = {
 1: '[1] M. Xhemrishi, J. \u00d6stman, A. Wachter-Zeh, and A. Graell i Amat, "FedGT: Identification of '
    'malicious clients in federated learning with secure aggregation," IEEE Trans. Inf. Forensics '
    'Security, vol. 20, pp. 2577\u20132592, 2025.',
 2: '[2] K. Bonawitz, V. Ivanov, B. Kreuter, A. Marcedone, H. B. McMahan, S. Patel, D. Ramage, A. Segal, '
    'and K. Seth, "Practical secure aggregation for privacy-preserving machine learning," in Proc. ACM '
    'CCS, 2017, pp. 1175\u20131191.',
 3: '[3] National Institute of Standards and Technology, "Status report on the fourth round of the NIST '
    'post-quantum cryptography standardization process," NIST IR 8545, Mar. 2025.',
 4: '[4] P. W. Shor, "Polynomial-time algorithms for prime factorization and discrete logarithms on a '
    'quantum computer," SIAM J. Comput., vol. 26, no. 5, pp. 1484\u20131509, 1997.',
 5: '[5] X. Cao, M. Fang, J. Liu, and N. Z. Gong, "FLTrust: Byzantine-robust federated learning via trust '
    'bootstrapping," in Proc. NDSS, 2021.',
 6: '[6] D. Yin, Y. Chen, R. Kannan, and P. Bartlett, "Byzantine-robust distributed learning: Towards '
    'optimal statistical rates," in Proc. ICML, 2018, pp. 5650\u20135659.',
 7: '[7] P. Blanchard, E. M. El Mhamdi, R. Guerraoui, and J. Stainer, "Machine learning with adversaries: '
    'Byzantine tolerant gradient descent," in Proc. NeurIPS, 2017, pp. 119\u2013129.',
 8: '[8] K. Pillutla, S. M. Kakade, and Z. Harchaoui, "Robust aggregation for federated learning," IEEE '
    'Trans. Signal Process., vol. 70, pp. 1142\u20131154, 2022.',
 9: '[9] National Institute of Standards and Technology, "Module-lattice-based key-encapsulation mechanism '
    'standard," FIPS 203, Aug. 2024.',
}


def tb(s, l, t, w, h):
    b = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    b.text_frame.word_wrap = True
    return b


def run(p, text, size, color, bold=False, italic=False):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.name = BODY; r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return r


def run_md(p, text, size, color):
    """Render a bullet; **phrase** is emphasized (bold + ink)."""
    parts = text.split("**")
    for i, seg in enumerate(parts):
        if seg == "":
            continue
        run(p, seg, size, INK if i % 2 else color, bold=(i % 2 == 1))


def footer(s, page):
    run(tb(s, 0.75, 7.05, 8.4, 0.3).text_frame.paragraphs[0],
        "CB-SAFE \u00b7 Suva & Chowdhury \u00b7 AIUB", 10.5, MUT)
    if page:
        pn = tb(s, 12.05, 7.05, 1.0, 0.3); pp = pn.text_frame.paragraphs[0]
        pp.alignment = PP_ALIGN.RIGHT; run(pp, page, 10.5, MUT)


def cites(s, keys):
    if not keys:
        return
    bx = tb(s, 0.75, 6.42, 11.9, 0.58).text_frame
    for i, k in enumerate(keys):
        p = bx.paragraphs[0] if i == 0 else bx.add_paragraph()
        p.space_after = Pt(1)
        run(p, REF[k], 8, BLUE, italic=True)


def callout(s, l, t, w, h, lines):
    box = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = CALL
    box.line.color.rgb = RGBColor(0xD6, 0xE0, 0xF0); box.line.width = Pt(0.75)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.12); tf.margin_bottom = Inches(0.1)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4); run_md(p, ln, 13, INK2)
    return box


def img_hw(name):
    w, h = Image.open(FIG + name).size
    return w / h


def content(kicker, title, bullets, cite=None, page=None, lead=None,
            fig=None, fig_center=None, bullet_w=None, bold=False):
    s = P.slides.add_slide(BLANK)
    run(tb(s, 0.75, 0.42, 11.9, 0.35).text_frame.paragraphs[0], kicker, 12, BLUE, bold=True)
    run(tb(s, 0.75, 0.78, 11.9, 0.9).text_frame.paragraphs[0], title, 27, INK, bold=True)
    top = 1.72
    if lead:
        run(tb(s, 0.75, top, 11.9, 0.5).text_frame.paragraphs[0], lead, 14, INK2, italic=True)
        top += 0.52
    bw = bullet_w if bullet_w else (7.0 if fig else 11.9)
    if bullets:
        body = tb(s, 0.75, top, bw, 4.3).text_frame
        for i, b in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.space_after = Pt(8)
            r0 = p.add_run(); r0.text = "\u25aa  "
            r0.font.size = Pt(15); r0.font.name = BODY; r0.font.color.rgb = BLUE
            (run_md if bold else run)(p, b, 15, INK2) if bold else run(p, b, 15, INK2)
    if fig:
        ar = img_hw(fig)
        W = 4.9; H = W / ar
        s.shapes.add_picture(FIG + fig, Inches(8.0), Inches(top + 0.15), width=Inches(W))
    if fig_center:
        name, ftop, fh = fig_center
        ar = img_hw(name)
        W = fh * ar
        if W > 12.4:
            W = 12.4; fh = W / ar
        s.shapes.add_picture(FIG + name, Inches((13.333 - W) / 2), Inches(ftop), height=Inches(fh))
    cites(s, cite or [])
    footer(s, page)
    return s


# ===================== TITLE =====================
s = P.slides.add_slide(BLANK)
run(tb(s, 0.75, 1.95, 11.9, 0.4).text_frame.paragraphs[0],
    "POST-QUANTUM CRYPTOGRAPHY \u00d7 FEDERATED LEARNING \u00b7 METHOD & FULL-PAPER TALK", 13, BLUE, bold=True)
run(tb(s, 0.75, 2.5, 11.7, 2.0).text_frame.paragraphs[0],
    "CB-SAFE: Code-Based Post-Quantum Secure Aggregation for Byzantine-Robust Federated Learning",
    37, INK, bold=True)
run(tb(s, 0.75, 4.95, 11.7, 0.5).text_frame.paragraphs[0],
    "Md Wahiduzzaman Suva (26-94088-2)   \u00b7   Esm-e Moula Chowdhury Abha (26-94089-2)", 17, INK2)
run(tb(s, 0.75, 5.42, 11.7, 0.4).text_frame.paragraphs[0],
    "Department of Computer Science, American International University\u2013Bangladesh", 14, MUT)
run(tb(s, 0.75, 6.35, 11.7, 0.4).text_frame.paragraphs[0],
    "Slides 1\u201315: required method presentation.   Slides 16\u201322: full-paper research walkthrough.", 12.5, MUT, italic=True)

# ===================== PART I: METHOD PRESENTATION (1-15) =====================
content("PROBLEM STATEMENT", "Three threats collide in one protocol", [
    "Federated learning shares model updates, not raw data, yet **those updates still leak private training data**.",
    "Secure aggregation hides updates behind masks, but its key exchange **breaks under a quantum computer** (harvest-now, decrypt-later).",
    "Every post-quantum fix today is **lattice-only, a monoculture**; and hiding updates also blinds the server to poisoned ones.",
    "Goal: one protocol that is **private, post-quantum with crypto diversity, and Byzantine-robust**.",
], cite=[2], page="1 / 22", bold=True)

content("BACKGROUND / MOTIVATION", "Why the problem is hard and timely", [
    "Secure aggregation = pairwise masks that cancel in the sum; the server sees only the aggregate, never an individual.",
    "Quantum threat: Shor's algorithm breaks classical key exchange, and traffic recorded today is decryptable later.",
    "Crypto diversity: NIST selected the code-based KEM HQC in 2025 to hedge a lattice break; that diversity has not reached FL.",
    "Robustness tension: privacy hides exactly the per-client signal that poisoning defenses need.",
], cite=[3, 4], page="2 / 22")

content("RELATED WORK", "FedGT: malicious-client ID under secure aggregation", [
    "FedGT (IEEE TIFS, 2025) is the closest peer to our robustness mechanism.",
    "Idea: group clients into overlapping test groups; the server sees only group aggregates, so privacy is preserved.",
    "A statistical test flags groups that contain malicious clients; a decoder then identifies and removes the culprits.",
    "Borrowed from group testing: find a few defectives with far fewer pooled tests than testing everyone.",
], cite=[1], page="3 / 22")

content("EXISTING METHOD", "How FedGT identifies attackers", [
    "An assignment matrix groups n clients into m overlapping groups (the parity-check structure of a code).",
    "Each group is securely aggregated; a test on the group aggregate is positive if any member is malicious.",
    "A COMP / Neyman-Pearson decoder infers which clients are malicious from the group test results.",
    "Flagged clients are excluded; the remaining clean clients are aggregated into the global model.",
], cite=[1], page="4 / 22")

content("LIMITATIONS", "Where FedGT leaves gaps (our angle)", [
    "Crypto: it treats secure aggregation as a black box, with **no post-quantum layer**, in a lattice-only field.",
    "Cost: it needs a **server validation set** and per-round group testing over fixed overlapping groups.",
    "Single-round decode: it identifies attackers from **one round's tests**, using no information across rounds.",
    "Unexplained failure: it does not characterize **why ordinary robust aggregation breaks** once updates are hidden.",
], cite=[1], page="5 / 22", bold=True)

content("WHY THIS PAPER", "Why FedGT is our basis (selection criteria)", [
    "Venue quality: IEEE TIFS, a Q1 security journal. Recency: published 2025.",
    "Relevance: it solves our exact sub-problem, malicious-client identification under secure aggregation.",
    "Method strength: its group testing outperforms the geometric median (RFA) and Multi-Krum in its own study.",
    "Reproducible: public code and a clearly specified protocol, a solid foundation to build on.",
], cite=[1, 8], page="6 / 22")

content("OUR PROPOSAL", "CB-SAFE: code-based, private, and robust", [
    "Hide updates within clusters using a code-based (HQC) masking protocol: the server sees only cluster sums.",
    "Run group testing across re-randomized clusters over rounds to identify attackers (CB-SAFE+).",
    "Confidentiality lives below the cluster boundary; robustness operates above it; post-quantum by construction.",
], lead="One framework that is private, post-quantum, and Byzantine-robust at the same time.",
    cite=[1, 2], page="7 / 22")

content("COMPONENT 1 / 3", "Taken from secure aggregation: masking", [
    "We adopt Bonawitz-style double masking: pairwise masks cancel in the cluster sum, so individuals stay hidden.",
    "Dropout resilience via Shamir secret sharing of the self-masks.",
    "Our change: masks are seeded from a KEM, so per-round traffic carries no cryptographic bytes at all.",
], cite=[2], page="8 / 22", fig="fig_masking.png")

content("COMPONENT 2 / 3", "Taken from FedGT: group testing for detection", [
    "We adopt the core idea: identify malicious clients from group-level observations, never from individual updates.",
    "The groups are the privacy-preserving cluster sums (the pools); a probe flags contaminated pools.",
    "Our change: tests repeat over re-randomized clusters and accumulate across rounds (temporal group testing).",
], cite=[1], page="9 / 22")

content("COMPONENT 3 / 3", "Taken from NIST PQC: a code-based KEM (HQC)", [
    "We adopt HQC, the code-based KEM NIST selected in 2025, for post-quantum confidentiality.",
    "Its security rests on syndrome decoding of quasi-cyclic codes, not lattices: genuine cryptographic diversity.",
    "Our change: HQC sits behind a KEM-agnostic interface, so ML-KEM swaps in with one configuration line.",
], cite=[3, 9], page="10 / 22")

content("COMBINATION", "How the three ideas fit together", [
    "HQC-seeded masking [2], [3] hides updates inside clusters, so the server sees only cluster sums.",
    "Group testing [1] then runs across those cluster sums, over re-randomized rounds, to name the attackers.",
    "Net design: privacy below the cluster boundary, robustness above it, code-based crypto throughout.",
    "On top we add our own pieces: the laundering theorem and a server root-loss probe.",
], lead="Masking + group testing + a code-based KEM become one coherent framework.",
    cite=[1, 2], page="11 / 22")

content("COMPARISON", "Previous (FedGT) vs proposed (CB-SAFE)", [
    "Crypto: FedGT is crypto-agnostic \u2192 **CB-SAFE is code-based post-quantum**.",
    "Detection: FedGT uses single-round COMP + a validation set \u2192 **CB-SAFE uses temporal accumulation + a root-loss probe** (plus a trust-free variant).",
    "Theory: FedGT gives no failure model \u2192 **CB-SAFE proves the laundering boundary**.",
    "Results: accuracy ties FedGT with **fewer false exclusions and ~4\u00d7 lower cost**.",
], cite=[1], page="12 / 22", fig="fig_signflip_acc_slide.png", bold=True)

content("WHAT IS NEW", "Our genuinely new contributions", [
    "**First code-based post-quantum secure aggregation** for FL: cryptographic diversity beyond lattices.",
    "The **laundering theorem**: at \u03b3*=2c/m\u22121 a poisoned cluster mean is magnitude-invisible, which explains why robust rules fail.",
    "**Temporal group testing** over re-randomized clusters with a root-loss probe: identifies attackers the server never sees individually.",
    "A **trust-free variant** that needs no root dataset, a capability FedGT does not offer.",
], cite=[1, 5], page="13 / 22", fig="fig_launder_geo.png", bold=True)

content("PARAMETERS", "How the selection meets the required criteria", [
    "Venue and recency: the basis paper is Q1 (IEEE TIFS), 2025; supporting sources are top-tier (ACM CCS, NIST).",
    "Relevance: every borrowed idea maps directly onto a component of our exact problem.",
    "Method strength: we match SOTA (FedGT) on robustness at ~4\u00d7 lower cost, add PQC, and prove the failure mechanism.",
    "Breadth and rigor: validated on four datasets and two modalities, multiple seeds, with significance testing.",
], cite=[1, 2, 3], page="14 / 22")

content("SUMMARY", "What we took, changed, and made new", [
    "Took: masking [2], group testing [1], and a code-based KEM [3].",
    "Changed: temporal group testing over re-randomized clusters; HQC-seeded, KEM-agnostic masking.",
    "New: the laundering theorem, the first code-based PQ secure aggregation, and trust-free detection.",
    "Payoff: private, post-quantum, and Byzantine-robust FL that ties the state of the art at lower cost.",
], lead="End of the required method presentation. Slides 16\u201322 present the full paper and results.",
    cite=[1, 2, 3], page="15 / 22")

# ===================== PART II: FULL-PAPER WALKTHROUGH (16-22) =====================
content("PART II \u00b7 SYSTEM ARCHITECTURE", "CB-SAFE: one training round, end to end", [],
    lead="Confidentiality within clusters, then temporal group testing over re-randomized clusters, then robust aggregation.",
    cite=[1, 2], page="16 / 22",
    fig_center=("fig_architecture.png", 2.55, 3.9))

content("EXPERIMENTAL SETUP", "Four datasets, two modalities, non-IID", [
    "Images: CIFAR-10, FashionMNIST, EMNIST-balanced (47 classes). Tabular: Edge-IIoTset (IoT intrusion).",
    "Non-IID split via a Dirichlet(0.5) partition; 30 communication rounds; clusters of size c = 3.",
    "Attacks: sign-flip (laundering), label-flip, and backdoor. Baselines: mean, trimmed, median, Krum, Bulyan, RFA, FLTrust, FedGT.",
], cite=[1, 6], page="17 / 22", fig="fig_dirichlet_slide.png")

content("MAIN RESULT (1 / 2) \u00b7 ROBUSTNESS", "Baselines collapse; CB-SAFE+ holds", [
    "Under sign-flip at f \u2265 0.2, **every coordinate-wise rule collapses to chance**; laundering makes poisoned means invisible.",
    "**CB-SAFE+ tracks the clean baseline** across all four datasets and **ties FedGT** while beating every other defense.",
], cite=[1], page="18 / 22",
    fig_center=("fig_dynamics_signflip_slide.png", 3.02, 3.28), bold=True)

content("MAIN RESULT (2 / 2) \u00b7 DETECTION", "It names attackers the server never sees", [
    "Malicious suspicion separates cleanly from honest across rounds, matching the predicted honest rate p_h.",
    "**Equal catch to FedGT (3/3, 6/6, 9/9) with far fewer false exclusions: 0.0 / 0.7 / 0.3 vs 1.3 / 3.3 / 1.3**.",
], cite=[1], page="19 / 22",
    fig_center=("fig_suspicion_slide.png", 3.2, 3.0), bold=True)

s20 = content("GENERALIZATION", "No robust rule collapses under label-flip", [
    "Label-flip is a milder attack, and it doubles as a control: unlike sign-flip, coordinate-wise rules do not collapse.",
    "That contrast isolates laundering as the mechanism behind the sign-flip failures.",
    "CB-SAFE+ stays at or above every baseline on all four datasets.",
], cite=[1], page="20 / 22", bullet_w=7.2)
callout(s20, 8.05, 2.05, 4.55, 3.2, [
    "**CB-SAFE+ label-flip accuracy (best per column):**",
    "EMNIST: **85.8 / 85.5 / 85.4** at f = .1 / .2 / .3",
    "FashionMNIST: **86.4 / 86.5 / 86.2**",
    "CIFAR-10: 53.6 / 53.2 / 51.1",
    "Edge-IIoTset: 62.4 / 62.5 / 62.4",
    "Undefended mean falls to 47\u201363% at f = .3, versus 10% under sign-flip.",
])

content("EFFICIENCY & THEORY", "Post-quantum diversity at lower cost", [
    "Code-based (HQC) confidentiality gives crypto diversity beyond the lattice monoculture; ML-KEM swaps in unchanged.",
    "One-time setup 15.2 KiB (HQC) vs 3.8 KiB (ML-KEM), paid once, under 0.03% of total traffic; per-round updates are 8d bytes regardless of KEM.",
    "Detection needs ~4\u00d7 fewer operations than FedGT (root-loss probe over cluster sums).",
    "Laundering theorem: at \u03b3* = 2c/m \u2212 1 the poisoned cluster mean is magnitude-invisible, which is why robust rules fail.",
], cite=[3, 9], page="21 / 22", fig="fig_launder_geo.png")

content("CONCLUSION", "Private, post-quantum, and Byzantine-robust", [
    "CB-SAFE unifies code-based PQ secure aggregation with temporal group-testing robustness in one protocol.",
    "It ties the state of the art (FedGT) on robustness and detection, at lower cost, and adds crypto diversity plus a failure theory.",
    "Limitations: update authentication is out of the code-based scope; EMNIST detection uses a single seed.",
    "Future work: code-based signatures for authentication, larger client populations, and adaptive attackers.",
], cite=[1, 2, 3], page="22 / 22")

# ===================== REFERENCES =====================
rs = P.slides.add_slide(BLANK)
run(tb(rs, 0.75, 0.42, 11.8, 0.35).text_frame.paragraphs[0], "REFERENCES", 12, BLUE, bold=True)
run(tb(rs, 0.75, 0.78, 11.9, 0.9).text_frame.paragraphs[0], "References", 27, INK, bold=True)
tf = tb(rs, 0.75, 1.75, 11.9, 5.1).text_frame
for i in range(1, 10):
    p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
    p.space_after = Pt(8); run(p, REF[i], 11.5, INK2)
run(tb(rs, 0.75, 7.05, 8.4, 0.3).text_frame.paragraphs[0],
    "CB-SAFE \u00b7 Suva & Chowdhury \u00b7 AIUB", 10.5, MUT)

out = r"d:/ISM_Quantum/paper/cbsafe_deck.pptx"
P.save(out)
print("saved", out, "| total slide objects:", len(P.slides._sldIdLst))
