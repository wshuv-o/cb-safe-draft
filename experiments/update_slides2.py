# -*- coding: utf-8 -*-
"""Remove the left blue accent bar from every slide; add the laundering
illustration to slide 7. Preserves everything else."""
import shutil
from pptx import Presentation
from pptx.util import Emu, Inches

PPTX = r"d:/ISM_Quantum/paper/cbsafe_slides.pptx"
FIG = r"d:/ISM_Quantum/results/figs/fig_launder_geo.png"

P = Presentation(PPTX)

# 1) remove the blue accent bar (the shape anchored at the very left edge)
removed = 0
for s in P.slides:
    for sh in list(s.shapes):
        if sh.left is not None and sh.top is not None:
            if abs(Emu(sh.left).inches) < 0.02 and abs(Emu(sh.top).inches) < 0.02 \
               and Emu(sh.width).inches < 0.25 and Emu(sh.height).inches > 5:
                sh._element.getparent().remove(sh._element)
                removed += 1
print("blue bars removed:", removed)

# 2) slide 7: shrink the lead-bullet box, add the laundering illustration top-right
s7 = P.slides[6]
for sh in s7.shapes:
    if sh.has_text_frame and "Key insight" in sh.text_frame.text:
        sh.width = Inches(8.7)          # was 11.90; free the right side
        break
s7.shapes.add_picture(FIG, Inches(9.75), Inches(1.80), Inches(2.67), Inches(2.00))
# small caption under the figure? place a tiny label
tb = s7.shapes.add_textbox(Inches(9.55), Inches(3.78), Inches(3.10), Inches(0.24))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Poison averages to normal size, wrong direction"
from pptx.dml.color import RGBColor
r.font.size = Emu(int(0.09 * 914400))   # ~9pt
r.font.name = "Segoe UI"
r.font.italic = True
r.font.color.rgb = RGBColor(0x89, 0x87, 0x81)

P.save(PPTX)
print("saved")
