# -*- coding: utf-8 -*-
"""Expand CB-SAFE deck from proposal (12) to a full-paper presentation:
15 content slides + title + references (17 total). Preserves the DNA."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPTX = r"d:/ISM_Quantum/paper/cbsafe_slides.pptx"
BLUE = RGBColor(0x2A, 0x78, 0xD6)
INK = RGBColor(0x0B, 0x0B, 0x0B)
INK2 = RGBColor(0x26, 0x25, 0x1F)
MUT = RGBColor(0x89, 0x87, 0x81)
BODY = "Segoe UI"

P = Presentation(PPTX)
LAYOUT = P.slides[1].slide_layout   # blank layout used by content slides


def _tb(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb


def _run(p, text, size, color, bold=False, italic=False):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.name = BODY; r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return r


def add_slide(kicker, title, bullets, lead=None, page=None):
    s = P.slides.add_slide(LAYOUT)
    # strip any inherited placeholders
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    _run(_tb(s, 0.75, 0.42, 11.8, 0.35).text_frame.paragraphs[0], kicker, 12, BLUE, bold=True)
    _run(_tb(s, 0.75, 0.75, 11.9, 0.95).text_frame.paragraphs[0], title, 30, INK, bold=True)
    top = 1.75
    if lead:
        _run(_tb(s, 0.75, top, 11.9, 0.5).text_frame.paragraphs[0], lead, 15, INK2)
        top = 2.35
    body = _tb(s, 0.75, top, 11.9, 4.6).text_frame
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.space_after = Pt(8)
        _run(p, "\u25aa  " + b, 15, INK2)
    # footers
    _run(_tb(s, 0.75, 7.05, 6.0, 0.3).text_frame.paragraphs[0],
         "CB-SAFE \u00b7 Suva & Chowdhury \u00b7 AIUB", 10.5, MUT)
    pn = _tb(s, 12.10, 7.05, 0.9, 0.3)
    pp = pn.text_frame.paragraphs[0]; pp.alignment = PP_ALIGN.RIGHT
    _run(pp, page or "", 10.5, MUT)
    return s


# ---- 4 new content slides (appended; reordered below) ----
add_slide("CONTRIBUTIONS", "What this paper contributes", [
    "First code-based (HQC) post-quantum secure aggregation for FL: cryptographic diversity beyond lattices.",
    "A quantified privacy\u2013robustness trade-off: clean-cluster probability (1\u2212f)^c ties the anonymity set to robustness.",
    "The laundering effect: a proved boundary \u03b3*=2c/m\u22121 where poisoning becomes magnitude-invisible.",
    "CB-SAFE+: temporal group testing that identifies attackers from group-level observations, at zero extra leakage.",
    "Evaluation on four datasets and two modalities; ties the SOTA (FedGT) at lower cost. Open-source.",
])
add_slide("GUARANTEE \u00b7 PRIVACY", "Provable confidentiality of masked aggregation", [
    "Theorem 1: an honest-but-curious server, even colluding with up to t\u22121 members of a cluster, learns nothing beyond each cluster's sum.",
    "Proof idea: IND-CCA2 KEM gives uniform pairwise secrets; the PRG/KDF give uniform masks; Shamir shares below threshold reveal nothing.",
    "Anonymity set = cluster size c. Re-randomizing clusters every round preserves the per-round guarantee.",
    "CB-SAFE+ flags are computed from sums the server already holds, so the defense adds no leakage.",
], lead="Leakage is exactly the per-cluster sum, and nothing else.")
add_slide("KEY FINDING \u00b7 LAUNDERING", "Secure aggregation can launder poisoning", [
    "Proposition 1: with m amplified sign-flips in a cluster of size c, the cluster mean is \u03bc = (c\u2212m(1+\u03b3))/c \u00b7 h.",
    "At \u03b3*=2c/m\u22121 the poisoned mean has the SAME magnitude as an honest update but opposite direction.",
    "So magnitude- and order-statistic rules (trimmed mean, median, even RFA) cannot see it and collapse to chance.",
    "Detection must therefore use signals that survive averaging: direction and repetition over rounds.",
])
add_slide("RESULTS \u00b7 GENERALIZATION", "Same story across domains; the privacy dial", [
    "Cross-domain: the collapse-and-recover pattern replicates on EMNIST (85% at f=30%) and tabular Edge-IIoTset (62%).",
    "Cluster-size dial: larger c enlarges the anonymity set but shifts the median's breakdown frontier left, exactly as (1\u2212f)^c predicts.",
    "Overlap helps: overlapping re-randomized clusters extend CB-SAFE+ recovery to higher f.",
    "Statistical rigor: 3 seeds on CIFAR-10/FashionMNIST/Edge-IIoTset, paired t-test with Holm correction.",
])

# ---- References slide (two columns) ----
refs_left = [
    "[1] McMahan et al., Communication-efficient learning, AISTATS 2017.",
    "[2] Bonawitz et al., Practical secure aggregation, CCS 2017.",
    "[3] Bell et al., Secure single-server aggregation, CCS 2020.",
    "[4] Xhemrishi et al., FedGT, IEEE TIFS 2025.",
    "[5] Cao et al., FLTrust, NDSS 2021.",
    "[6] Blanchard et al., Krum, NeurIPS 2017.",
    "[7] Yin et al., Byzantine-robust distributed learning, ICML 2018.",
    "[8] El Mhamdi et al., Bulyan, ICML 2018.",
]
refs_right = [
    "[9] Pillutla et al., RFA (geometric median), IEEE TSP 2022.",
    "[10] Shor, Quantum factoring, SIAM J. Comput. 1997.",
    "[11] McEliece, Code-based cryptosystem, JPL DSN 1978.",
    "[12] NIST IR 8545, PQC round-4 status, 2025 (HQC).",
    "[13] NIST FIPS 203, ML-KEM, 2024.",
    "[14] Bagdasaryan et al., Backdoor FL, AISTATS 2020.",
    "[15] Ferrag et al., Edge-IIoTset, IEEE Access 2022.",
    "Full list (47 refs) in the manuscript.",
]
rs = P.slides.add_slide(LAYOUT)
for ph in list(rs.placeholders):
    ph._element.getparent().remove(ph._element)
_run(_tb(rs, 0.75, 0.42, 11.8, 0.35).text_frame.paragraphs[0], "REFERENCES", 12, BLUE, bold=True)
_run(_tb(rs, 0.75, 0.75, 11.9, 0.95).text_frame.paragraphs[0], "References", 30, INK, bold=True)
for col, items in ((0.75, refs_left), (6.9, refs_right)):
    tf = _tb(rs, col, 1.9, 6.0, 5.0).text_frame
    for i, r in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        _run(p, r, 11, INK2)
_run(_tb(rs, 0.75, 7.05, 6.0, 0.3).text_frame.paragraphs[0],
     "CB-SAFE \u00b7 Suva & Chowdhury \u00b7 AIUB", 10.5, MUT)

# ---- reframe slide 5 kicker (proposal -> approach) ----
for sh in P.slides[4].shapes:
    if sh.has_text_frame and sh.text_frame.text.strip() == "OUR PROPOSAL":
        sh.text_frame.paragraphs[0].runs[0].text = "APPROACH OVERVIEW"

# ---- reorder sldIdLst to the full-paper sequence ----
# current 0-based: 0 title,1 intro,2 bg,3 rel,4 approach,5 conf,6 cbsafe+,7 method,
#                  8 cost,9 util,10 robust,11 concl, 12 contrib,13 privacy,14 launder,
#                  15 general,16 refs
order = [0, 1, 2, 3, 12, 4, 5, 13, 14, 6, 7, 8, 9, 10, 15, 11, 16]
sldIdLst = P.slides._sldIdLst
ids = list(sldIdLst)
for el in ids:
    sldIdLst.remove(el)
for idx in order:
    sldIdLst.append(ids[idx])

# ---- renumber page numbers: content 1..15, title & refs blank ----
# after reorder, slide positions: 0 title, 1..15 content, 16 refs
slides = list(P.slides)
for pos, s in enumerate(slides):
    num = "" if (pos == 0 or pos == 16) else f"{pos} / 15"
    for sh in s.shapes:
        if sh.has_text_frame and "/" in sh.text_frame.text and sh.left is not None \
           and Emu(sh.left).inches > 11.5:
            sh.text_frame.paragraphs[0].runs[0].text = num
        elif sh.has_text_frame and sh.text_frame.text.strip() in ("", None) and False:
            pass
    # title/refs may lack a page-number box; content slides have one from add_slide/original

P.save(PPTX)
print("expanded to", len(list(P.slides)), "slides")
