# -*- coding: utf-8 -*-
"""Update cbsafe_slides.pptx in place with current results, preserving the DNA
(blue accent bar, Segoe UI, bullet callouts). Only content is changed."""
from pptx import Presentation

PPTX = r"d:/ISM_Quantum/paper/cbsafe_slides.pptx"
FIG = r"d:/ISM_Quantum/results/figs/fig_signflip_acc_slide.png"


def find(slide, contains):
    for sh in slide.shapes:
        if sh.has_text_frame and contains in sh.text_frame.text:
            return sh
    return None


def set_paras(shape, texts):
    """Set paragraph texts in place, preserving each paragraph's first-run format."""
    paras = shape.text_frame.paragraphs
    assert len(paras) >= len(texts), f"need {len(texts)} paras, have {len(paras)}"
    for i, t in enumerate(texts):
        p = paras[i]
        if p.runs:
            p.runs[0].text = t
            for r in p.runs[1:]:
                r._r.getparent().remove(r._r)
        else:
            p.add_run().text = t


P = Presentation(PPTX)
S = P.slides

# --- Slide 4: related work callout -> mention FedGT ---
sh = find(S[3], "CB-SAFE targets exactly these gaps")
set_paras(sh, [
    "CB-SAFE targets exactly these gaps. The closest robustness peer, FedGT "
    "(TIFS 2025), also tests group sums but is not post-quantum; CB-SAFE matches "
    "its detection at lower cost and adds code-based confidentiality."
])

# --- Slide 8: methodology baselines -> full panel ---
sh = find(S[7], "Baselines: plain FedAvg")
set_paras(sh, [
    "Baselines: plain FedAvg and the lattice (ML-KEM) variant, plus eight robust "
    "rules at matched NIST levels: mean, trimmed mean, median, Multi-Krum, Bulyan, "
    "geometric-median (RFA), FLTrust, and the SOTA group-testing method FedGT "
    "(TIFS 2025)."
])

# --- Slide 11: robustness bullets + caption + figure swap ---
sh = find(S[10], "Why rules fail")
set_paras(sh, [
    "\u25aa  Why rules fail: cluster averaging disguises the poison (normal size, "
    "wrong direction), so mean, trimmed mean, median and even the robust geometric "
    "median (RFA) fall to ~10% (random) at f \u2265 20%; Multi-Krum and Bulyan "
    "degrade but hold longer.",
    "\u25aa  Ties the state of the art: at f = 30% CB-SAFE+ holds ~50% accuracy, "
    "matching FedGT (TIFS 2025), while catching 3/3, 6/6, 9/9 attackers with far "
    "fewer false exclusions (0.0 / 0.7 / 0.3 vs FedGT 1.3 / 3.3 / 1.3).",
    "\u25aa  Generalizes across domains: the same collapse-and-recover pattern "
    "holds on EMNIST (85% at f = 30%) and tabular Edge-IIoTset (62%); stealthy "
    "backdoors remain an open limit for all sum-level defenses.",
])
cap = find(S[10], "Sign-flip attack")
set_paras(cap, [
    "Sign-flip attack on CIFAR-10: accuracy vs attacker fraction across eight "
    "rules; only CB-SAFE+ and FedGT stay near the clean baseline."
])
# swap the image (keep geometry)
pic = next((sh for sh in S[10].shapes if sh.shape_type == 13), None)
L, T, W, H = pic.left, pic.top, pic.width, pic.height
pic._element.getparent().remove(pic._element)
S[10].shapes.add_picture(FIG, L, T, W, H)

# --- Slide 12: conclusion bullets ---
sh = find(S[11], "Diversity delivered")
set_paras(sh, [
    "\u25aa  Diversity delivered: the first code-based post-quantum secure "
    "aggregation for FL, at a one-time cost under 0.03% of traffic.",
    "\u25aa  New understanding: secure aggregation launders poisoning; defenses "
    "must use signals that survive summation, accumulated over rounds.",
    "\u25aa  New defense: CB-SAFE+ identifies individual attackers from group "
    "observations only, tying FedGT at ~4x lower cost with zero extra privacy leakage.",
    "\u25aa  Open problems: stealthy backdoors, actively malicious servers, "
    "adaptive attackers.",
    "\u25aa  Validated across four datasets (CIFAR-10, FashionMNIST, EMNIST, "
    "Edge-IIoTset); 100-client scale in progress; implementation released open-source.",
])

P.save(PPTX)
print("updated", PPTX, "| slides:", len(S))
