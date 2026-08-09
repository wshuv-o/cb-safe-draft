# -*- coding: utf-8 -*-
"""Build the instructor-structured 'method presentation' deck for CB-SAFE:
Title + Outline + 15 content slides + References, with per-slide citations.
Keeps the DNA (Segoe UI, blue kicker, no accent bar, muted footers)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLUE = RGBColor(0x2A, 0x78, 0xD6)
INK = RGBColor(0x0B, 0x0B, 0x0B)
INK2 = RGBColor(0x26, 0x25, 0x1F)
MUT = RGBColor(0x89, 0x87, 0x81)
BODY = "Segoe UI"
FIG = r"d:/ISM_Quantum/results/figs/"

P = Presentation()
P.slide_width = Inches(13.333)
P.slide_height = Inches(7.5)
BLANK = P.slide_layouts[6]

TOTAL = 15  # content slides


def tb(s, l, t, w, h):
    b = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    b.text_frame.word_wrap = True
    return b


def run(p, text, size, color, bold=False, italic=False):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.name = BODY; r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return r


def footer(s, page):
    run(tb(s, 0.75, 7.05, 8.0, 0.3).text_frame.paragraphs[0],
        "CB-SAFE \u00b7 Suva & Chowdhury \u00b7 AIUB", 10.5, MUT)
    if page:
        pn = tb(s, 12.10, 7.05, 0.9, 0.3); pp = pn.text_frame.paragraphs[0]
        pp.alignment = PP_ALIGN.RIGHT; run(pp, page, 10.5, MUT)


def content(kicker, title, bullets, cite=None, page=None, lead=None,
            fig=None, fig_wide=None):
    s = P.slides.add_slide(BLANK)
    run(tb(s, 0.75, 0.42, 11.8, 0.35).text_frame.paragraphs[0], kicker, 12, BLUE, bold=True)
    run(tb(s, 0.75, 0.78, 11.9, 0.9).text_frame.paragraphs[0], title, 28, INK, bold=True)
    top = 1.75
    if lead:
        run(tb(s, 0.75, top, 11.9, 0.5).text_frame.paragraphs[0], lead, 14, INK2, italic=True)
        top += 0.55
    bw = 7.0 if fig else 11.9
    body = tb(s, 0.75, top, bw, 4.6).text_frame
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.space_after = Pt(9)
        run(p, "\u25aa  " + b, 15, INK2)
    if fig:
        s.shapes.add_picture(FIG + fig, Inches(8.05), Inches(top + 0.05),
                             height=Inches(3.4) if False else None, width=Inches(4.6))
    if fig_wide:
        from PIL import Image
        w, h = Image.open(FIG + fig_wide).size
        H = 3.0
        W = H * (w / h)
        s.shapes.add_picture(FIG + fig_wide, Inches((13.333 - W) / 2), Inches(3.45), height=Inches(H))
    if cite:
        run(tb(s, 0.75, 6.62, 11.5, 0.32).text_frame.paragraphs[0], cite, 10, BLUE, italic=True)
    footer(s, page)
    return s


# ---------- Title ----------
s = P.slides.add_slide(BLANK)
run(tb(s, 0.75, 2.05, 11.8, 0.4).text_frame.paragraphs[0],
    "METHOD PRESENTATION \u00b7 POST-QUANTUM CRYPTOGRAPHY \u00d7 FEDERATED LEARNING", 13, BLUE, bold=True)
run(tb(s, 0.75, 2.55, 11.6, 2.0).text_frame.paragraphs[0],
    "CB-SAFE: Code-Based Post-Quantum Secure Aggregation for Byzantine-Robust Federated Learning",
    38, INK, bold=True)
run(tb(s, 0.75, 4.95, 11.6, 0.5).text_frame.paragraphs[0],
    "Md Wahiduzzaman Suva (26-94088-2)   \u00b7   Esm-e Moula Chowdhury Abha (26-94089-2)", 17, INK2)
run(tb(s, 0.75, 5.45, 11.6, 0.4).text_frame.paragraphs[0],
    "Department of Computer Science, American International University\u2013Bangladesh", 14, MUT)

# ---------- Outline / ToC ----------
content("OUTLINE", "What this talk covers", [
    "Problem statement, background and motivation (1\u20132).",
    "Closest related work (FedGT), its method, limitations, and why we selected it (3\u20136).",
    "Our proposed idea, and the component we take from each source paper (7\u201310).",
    "How we combine them, previous-vs-proposed, what is genuinely new (11\u201313).",
    "How the selection meets the required parameters, and the key contribution (14\u201315).",
], page=None)

# ---------- 1. Problem Statement ----------
content("PROBLEM STATEMENT", "Three threats collide in one protocol", [
    "Federated learning shares model updates, not raw data; but those updates leak private training data.",
    "Secure aggregation hides updates behind masks, yet its key exchange breaks under a quantum computer (harvest-now, decrypt-later).",
    "Every post-quantum fix today is lattice-only (a monoculture); and hiding updates also blinds the server to poisoned ones.",
    "Goal: one protocol that is private, post-quantum with crypto diversity, AND Byzantine-robust.",
], cite="[2] Bonawitz et al., Practical Secure Aggregation, ACM CCS 2017.", page="1 / 15")

# ---------- 2. Background / Motivation ----------
content("BACKGROUND / MOTIVATION", "Why the problem is hard and timely", [
    "Secure aggregation = pairwise masks that cancel in the sum; the server sees only the aggregate, never an individual.",
    "Quantum threat: Shor's algorithm breaks classical key exchange, and traffic recorded today is decryptable later.",
    "Crypto diversity: NIST selected the code-based KEM HQC in 2025 to hedge a lattice break; that diversity has not reached FL.",
    "Robustness tension: privacy hides exactly the per-client signal that poisoning defenses need.",
], cite="[3] NIST IR 8545, 2025 (HQC);  [4] Shor, SIAM J. Comput., 1997.", page="2 / 15")

# ---------- 3. Closest Related Work ----------
content("CLOSEST RELATED WORK", "FedGT: malicious-client ID under secure aggregation", [
    "FedGT (IEEE TIFS 2025) is the closest peer to our robustness mechanism.",
    "Idea: group clients into overlapping test groups; the server sees only group aggregates, so privacy is preserved.",
    "A statistical test flags groups that contain malicious clients; a decoder then identifies and removes the culprits.",
    "It borrows from group testing: find a few defectives with far fewer pooled tests than testing everyone.",
], cite="[1] Xhemrishi, \u00d6stman, Wachter-Zeh, Graell i Amat, FedGT, IEEE TIFS, vol. 20, 2025.", page="3 / 15")

# ---------- 4. Existing Method ----------
content("EXISTING METHOD", "How FedGT identifies attackers", [
    "An assignment matrix A groups n clients into m overlapping groups (the parity-check structure of a code).",
    "Each group is securely aggregated; a test on the group aggregate is positive if any member is malicious (a syndrome bit).",
    "A COMP / Neyman-Pearson decoder infers the defective vector (which clients are malicious) from the group test results.",
    "Flagged clients are excluded; the remaining clean clients are aggregated into the global model.",
], cite="[1] Xhemrishi et al., FedGT, IEEE TIFS 2025.", page="4 / 15")

# ---------- 5. Limitations ----------
content("LIMITATIONS", "Where FedGT leaves gaps (our angle)", [
    "Crypto: it treats secure aggregation as a black box, with no post-quantum layer, in a lattice-only field.",
    "Cost: it needs a server validation set and per-round group testing over fixed overlapping groups.",
    "Single-round decode: it identifies attackers from one round's tests, using no information across rounds.",
    "Unexplained failure: it does not characterize WHY ordinary robust aggregation breaks once updates are hidden.",
], cite="[1] Xhemrishi et al., FedGT, IEEE TIFS 2025.", page="5 / 15")

# ---------- 6. Why this paper ----------
content("WHY THIS PAPER", "Why FedGT is our basis (selection parameters)", [
    "Venue quality: IEEE TIFS, a Q1 security journal.  Recency: published 2025.",
    "Relevance: it solves our exact sub-problem, malicious-client identification under secure aggregation.",
    "Method strength: its group testing outperforms the geometric median (RFA) and Multi-Krum in its own study.",
    "Reproducible: public code and a clearly specified protocol, giving a solid, buildable foundation.",
], cite="[1] Xhemrishi et al., FedGT, IEEE TIFS 2025.", page="6 / 15")

# ---------- 7. Proposed Idea ----------
content("OUR PROPOSAL", "CB-SAFE: code-based, private, and robust", [
    "Hide updates within clusters using a code-based (HQC) masking protocol: the server sees only cluster sums.",
    "Run group testing ACROSS re-randomized clusters over rounds to identify attackers (CB-SAFE+).",
    "Confidentiality lives below the cluster boundary; robustness operates above it; post-quantum by construction.",
], cite="This work.", page="7 / 15", fig_wide="fig_arch_slide.png")

# ---------- 8. Component 1: secure aggregation ----------
content("COMPONENT 1 / 3", "Taken from secure aggregation: masking", [
    "We adopt Bonawitz-style double masking: pairwise masks cancel in the cluster sum, so individuals stay hidden.",
    "Dropout resilience via Shamir secret sharing of the self-masks.",
    "Our change: masks are seeded from a KEM, so per-round traffic carries no cryptographic bytes at all.",
], cite="[2] Bonawitz et al., Practical Secure Aggregation, ACM CCS 2017.", page="8 / 15", fig="fig_masking.png")

# ---------- 9. Component 2: group testing ----------
content("COMPONENT 2 / 3", "Taken from FedGT: group testing for detection", [
    "We adopt the core idea: identify malicious clients from GROUP-level observations, never from individual updates.",
    "The groups are the privacy-preserving cluster sums (the pools); a probe flags contaminated pools.",
    "Our change: tests repeat over re-randomized clusters and accumulate across rounds (temporal group testing), instead of a single-round decode.",
], cite="[1] Xhemrishi et al., FedGT, IEEE TIFS 2025.", page="9 / 15")

# ---------- 10. Component 3: code-based PQC ----------
content("COMPONENT 3 / 3", "Taken from NIST PQC: a code-based KEM (HQC)", [
    "We adopt HQC, the code-based KEM NIST selected in 2025, for post-quantum confidentiality.",
    "Its security rests on syndrome decoding of quasi-cyclic codes, not on lattices: genuine cryptographic diversity.",
    "Our change: HQC sits behind a KEM-agnostic interface, so one can swap in ML-KEM with a single configuration line.",
], cite="[3] NIST IR 8545, 2025 (HQC selection).", page="10 / 15")

# ---------- 11. How you combine ----------
content("COMBINATION", "How the three ideas fit together", [
    "HQC-seeded masking (from [2], [3]) hides updates inside clusters, so the server sees only cluster sums.",
    "Group testing (from [1]) then runs ACROSS those cluster sums, over re-randomized rounds, to name the attackers.",
    "Net design: privacy below the cluster boundary, robustness above it, code-based post-quantum crypto throughout.",
    "On top we add our own pieces (next slides): the laundering theorem and a server root-loss probe.",
], lead="Like combining a blue shirt and mustard pants: masking + group testing + code-based crypto become one framework.",
    cite="[1] FedGT 2025;  [2] Bonawitz 2017;  [3] NIST IR 8545 2025.", page="11 / 15")

# ---------- 12. Comparison ----------
content("COMPARISON", "Previous (FedGT) vs proposed (CB-SAFE)", [
    "Crypto: FedGT is crypto-agnostic (no PQC)  \u2192  CB-SAFE is code-based post-quantum.",
    "Detection: FedGT uses single-round COMP + a validation set  \u2192  CB-SAFE uses temporal accumulation + a root-loss probe (plus a trust-free variant with no root data).",
    "Theory: FedGT gives no failure model  \u2192  CB-SAFE proves the laundering boundary.",
    "Results: accuracy ties FedGT, with fewer false exclusions (0.0/0.7/0.3 vs 1.3/3.3/1.3) and ~4\u00d7 lower cost.",
], cite="[1] Xhemrishi et al., FedGT, IEEE TIFS 2025.", page="12 / 15", fig="fig_signflip_acc_slide.png")

# ---------- 13. What is new ----------
content("WHAT IS NEW", "Our genuinely new contributions", [
    "First code-based post-quantum secure aggregation for FL: cryptographic diversity beyond lattices.",
    "The laundering theorem: at \u03b3*=2c/m\u22121 a poisoned cluster mean is magnitude-invisible, which explains why robust rules fail.",
    "Temporal group testing over re-randomized clusters with a root-loss probe: identifies attackers the server never sees individually.",
    "A trust-free variant that needs no root dataset, a capability FedGT does not offer.",
], cite="This work;  root-loss probe inspired by [5] Cao et al., FLTrust, NDSS 2021.", page="13 / 15", fig="fig_launder_geo.png")

# ---------- 14. Parameters ----------
content("PARAMETERS", "How the selection meets the required criteria", [
    "Venue and recency: the basis paper FedGT is Q1 (IEEE TIFS), 2025; supporting sources are top-tier (ACM CCS, NIST).",
    "Relevance: every borrowed idea maps directly onto a component of our exact problem.",
    "Method strength: we match the SOTA (FedGT) on robustness at ~4\u00d7 lower cost, add PQC, and prove the failure mechanism.",
    "Breadth and rigor: validated on four datasets and two modalities, three seeds, with significance testing.",
], cite="[1] FedGT 2025;  [2] Bonawitz 2017;  [3] NIST IR 8545 2025.", page="14 / 15")

# ---------- 15. Summary ----------
content("SUMMARY", "What we took, changed, and made new", [
    "Took: masking [2], group testing [1], and a code-based KEM [3].",
    "Changed: temporal group testing over re-randomized clusters; HQC-seeded, KEM-agnostic masking.",
    "New: the laundering theorem, the first code-based PQ secure aggregation, and trust-free detection.",
    "Payoff: private, post-quantum, and Byzantine-robust FL that ties the state of the art at lower cost.",
], cite="[1] FedGT 2025;  [2] Bonawitz 2017;  [3] NIST IR 8545 2025.", page="15 / 15")

# ---------- References ----------
rs = P.slides.add_slide(BLANK)
run(tb(rs, 0.75, 0.42, 11.8, 0.35).text_frame.paragraphs[0], "REFERENCES", 12, BLUE, bold=True)
run(tb(rs, 0.75, 0.78, 11.9, 0.9).text_frame.paragraphs[0], "References", 28, INK, bold=True)
refs = [
    "[1] M. Xhemrishi, J. \u00d6stman, A. Wachter-Zeh, A. Graell i Amat. FedGT: Identification of malicious clients in federated learning with secure aggregation. IEEE Trans. Inf. Forensics Security, vol. 20, 2025.",
    "[2] K. Bonawitz et al. Practical secure aggregation for privacy-preserving machine learning. ACM CCS, 2017.",
    "[3] National Institute of Standards and Technology. Status report on the fourth round of the NIST PQC standardization process (HQC selection). NIST IR 8545, 2025.",
    "[4] P. W. Shor. Polynomial-time algorithms for prime factorization and discrete logarithms on a quantum computer. SIAM J. Comput., 1997.",
    "[5] X. Cao, M. Fang, J. Liu, N. Z. Gong. FLTrust: Byzantine-robust federated learning via trust bootstrapping. NDSS, 2021.",
]
tf = tb(rs, 0.75, 1.9, 11.9, 5.0).text_frame
for i, r in enumerate(refs):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(10); run(p, r, 12, INK2)
run(tb(rs, 0.75, 7.05, 8.0, 0.3).text_frame.paragraphs[0],
    "CB-SAFE \u00b7 Suva & Chowdhury \u00b7 AIUB", 10.5, MUT)

out = r"d:/ISM_Quantum/paper/cbsafe_method_slides.pptx"
P.save(out)
print("saved", out, "| slides:", len(P.slides.__iter__.__self__._sldIdLst))
